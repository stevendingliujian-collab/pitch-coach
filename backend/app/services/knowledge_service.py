"""
Knowledge base ingestion pipeline and hybrid search.

Ingestion flow:
  1. Download file from MinIO
  2. Parse document → raw text chunks (300-800 chars)
  3. Embed chunks via DashScope text-embedding-v3
  4. Upsert vectors into Qdrant
  5. Store chunk rows in PostgreSQL (knowledge_chunk table)
  6. Update knowledge_document.embedding_status + chunk_count

Hybrid search (RRF fusion):
  - Vector search: Qdrant cosine similarity → top-20
  - Keyword search: PostgreSQL tsvector → top-20
  - RRF(k=60) merge → reranked top-N
"""
from __future__ import annotations

import io
import re
import math
import asyncio
import logging
from typing import Any

from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select, text, update

from app.models.knowledge import KnowledgeDocument, KnowledgeChunk
from app.services.embedding_service import upsert_chunks, vector_search
from app.core.storage import get_minio_client
from app.core.config import get_settings

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Text chunking
# ---------------------------------------------------------------------------
MIN_CHUNK = 150
MAX_CHUNK = 800
OVERLAP = 50  # characters of overlap between consecutive chunks


def _split_text(text_content: str) -> list[str]:
    """
    Split text into overlapping chunks of 150-800 characters.
    Splits on paragraph boundaries (\n\n) first, then sentence boundaries (。！？).
    """
    # Normalise whitespace but keep paragraph breaks
    text_content = re.sub(r"\r\n", "\n", text_content)
    text_content = re.sub(r"[ \t]+", " ", text_content)

    # Split into paragraphs
    paragraphs = [p.strip() for p in re.split(r"\n{2,}", text_content) if p.strip()]

    chunks: list[str] = []
    current = ""

    for para in paragraphs:
        if len(current) + len(para) + 1 <= MAX_CHUNK:
            current = (current + "\n" + para).strip() if current else para
        else:
            # Flush current if long enough
            if len(current) >= MIN_CHUNK:
                chunks.append(current)
                # Overlap: keep last OVERLAP chars as seed
                current = current[-OVERLAP:] + "\n" + para
            else:
                # Current too short, just append
                current = (current + "\n" + para).strip()

            # If the new para itself exceeds MAX_CHUNK, split by sentences
            if len(current) > MAX_CHUNK:
                sentence_chunks = _split_by_sentence(current)
                chunks.extend(sentence_chunks[:-1])
                current = sentence_chunks[-1] if sentence_chunks else ""

    if current and len(current) >= MIN_CHUNK:
        chunks.append(current)
    elif current and chunks:
        # Append remainder to last chunk if too short
        chunks[-1] = chunks[-1] + " " + current

    return [c.strip() for c in chunks if c.strip()]


def _split_by_sentence(text_content: str) -> list[str]:
    """Split on Chinese/English sentence boundaries."""
    sentences = re.split(r"(?<=[。！？\.!?])\s*", text_content)
    chunks: list[str] = []
    current = ""
    for s in sentences:
        if len(current) + len(s) <= MAX_CHUNK:
            current += s
        else:
            if current:
                chunks.append(current.strip())
            current = s
    if current:
        chunks.append(current.strip())
    return [c for c in chunks if c.strip()]


# ---------------------------------------------------------------------------
# Document parsing
# ---------------------------------------------------------------------------

async def _extract_text(file_bytes: bytes, file_name: str) -> list[dict]:
    """
    Extract text from file bytes.
    Returns list of {"text": str, "page_number": int | None, "heading": str | None}.
    Supports: .txt, .md, .pdf, .docx, .pptx
    """
    name_lower = file_name.lower()

    if name_lower.endswith((".txt", ".md")):
        raw = file_bytes.decode("utf-8", errors="replace")
        return [{"text": raw, "page_number": None, "heading": None}]

    if name_lower.endswith(".pdf"):
        return await _extract_pdf(file_bytes)

    if name_lower.endswith(".docx"):
        return _extract_docx(file_bytes)

    if name_lower.endswith((".pptx", ".ppt")):
        return _extract_pptx(file_bytes)

    # Fallback: try UTF-8
    try:
        raw = file_bytes.decode("utf-8", errors="replace")
        return [{"text": raw, "page_number": None, "heading": None}]
    except Exception:
        return []


async def _extract_pdf(file_bytes: bytes) -> list[dict]:
    """Extract text from PDF using pdftotext (subprocess)."""
    import subprocess
    import tempfile
    import os

    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        result = subprocess.run(
            ["pdftotext", "-layout", tmp_path, "-"],
            capture_output=True, timeout=60
        )
        if result.returncode == 0:
            text_content = result.stdout.decode("utf-8", errors="replace")
            # Split by form-feed (page separator in pdftotext output)
            pages = text_content.split("\x0c")
            return [
                {"text": p.strip(), "page_number": i + 1, "heading": None}
                for i, p in enumerate(pages) if p.strip()
            ]
    except Exception as e:
        logger.warning(f"pdftotext failed: {e}")
    finally:
        os.unlink(tmp_path)

    return [{"text": "", "page_number": None, "heading": None}]


def _extract_docx(file_bytes: bytes) -> list[dict]:
    """Extract text from DOCX."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        paragraphs_by_heading: list[dict] = []
        current_heading: str | None = None
        current_text = ""

        for para in doc.paragraphs:
            if para.style.name.startswith("Heading"):
                if current_text.strip():
                    paragraphs_by_heading.append({
                        "text": current_text.strip(),
                        "page_number": None,
                        "heading": current_heading,
                    })
                current_heading = para.text.strip()
                current_text = ""
            else:
                current_text += para.text + "\n"

        if current_text.strip():
            paragraphs_by_heading.append({
                "text": current_text.strip(),
                "page_number": None,
                "heading": current_heading,
            })

        return paragraphs_by_heading or [
            {"text": "\n".join(p.text for p in doc.paragraphs), "page_number": None, "heading": None}
        ]
    except Exception as e:
        logger.warning(f"docx extraction failed: {e}")
        return []


def _extract_pptx(file_bytes: bytes) -> list[dict]:
    """Extract text from PPTX per slide."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        pages = []
        for i, slide in enumerate(prs.slides):
            texts = []
            title = None
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if shape.shape_type == 13:  # picture
                    continue
                t = shape.text_frame.text.strip()
                if t:
                    if shape.name.lower().startswith("title") or (
                        hasattr(slide.shapes, "title") and shape == slide.shapes.title
                    ):
                        title = t
                    else:
                        texts.append(t)
            full = "\n".join(texts)
            if full:
                pages.append({
                    "text": full,
                    "page_number": i + 1,
                    "heading": title,
                })
        return pages
    except Exception as e:
        logger.warning(f"pptx extraction failed: {e}")
        return []


# ---------------------------------------------------------------------------
# Main ingestion entry-point
# ---------------------------------------------------------------------------

async def ingest_document(
    doc_id: int,
    db: AsyncSession,
) -> None:
    """
    Full ingestion pipeline for a KnowledgeDocument.
    Called from Celery task (runs in asyncio context via asyncio.run).
    """
    # 1. Load document record
    doc: KnowledgeDocument = await db.get(KnowledgeDocument, doc_id)
    if doc is None:
        logger.error(f"KnowledgeDocument {doc_id} not found")
        return

    # Mark as processing
    doc.embedding_status = 1
    await db.commit()

    try:
        # 2. Download file from MinIO
        settings = get_settings()
        minio = get_minio_client()

        # file_url is like "pitchcoach/knowledge/filename.pdf"
        # strip leading bucket name if present
        object_key = doc.file_url
        if object_key.startswith(settings.minio_bucket + "/"):
            object_key = object_key[len(settings.minio_bucket) + 1:]

        response = minio.get_object(settings.minio_bucket, object_key)
        file_bytes = response.read()
        response.close()
        response.release_conn()

        # 3. Extract text
        pages = await _extract_text(file_bytes, doc.file_name)

        # 4. Chunk text
        raw_chunks: list[dict] = []
        for page_data in pages:
            page_text = page_data["text"]
            if not page_text.strip():
                continue
            text_chunks = _split_text(page_text)
            for chunk_text in text_chunks:
                raw_chunks.append({
                    "content": chunk_text,
                    "page_number": page_data.get("page_number"),
                    "heading": page_data.get("heading"),
                    "content_type": "text",
                })

        if not raw_chunks:
            logger.warning(f"No chunks extracted from document {doc_id}")
            doc.embedding_status = 3  # failed
            doc.chunk_count = 0
            await db.commit()
            return

        # 5. Insert chunk rows into PostgreSQL (to get chunk IDs)
        db_chunks: list[KnowledgeChunk] = []
        for idx, c in enumerate(raw_chunks):
            chunk = KnowledgeChunk(
                doc_id=doc_id,
                tenant_id=doc.tenant_id,
                chunk_index=idx,
                content=c["content"],
                content_type=c["content_type"],
                heading=c["heading"],
                page_number=c["page_number"],
            )
            db.add(chunk)
            db_chunks.append(chunk)

        await db.flush()  # get auto-generated IDs

        # 6. Embed + upsert to Qdrant
        embed_records = [
            {
                "chunk_id": chunk.id,
                "doc_id": doc_id,
                "tenant_id": doc.tenant_id,
                "content": chunk.content,
                "heading": chunk.heading,
                "page_number": chunk.page_number,
                "content_type": chunk.content_type,
            }
            for chunk in db_chunks
        ]
        point_ids = await upsert_chunks(embed_records)

        # 7. Write vector_id back to chunk rows
        for chunk, point_id in zip(db_chunks, point_ids):
            chunk.vector_id = point_id

        # 8. Update document status
        doc.chunk_count = len(db_chunks)
        doc.embedding_status = 2  # done

        await db.commit()
        logger.info(f"Document {doc_id} ingested: {len(db_chunks)} chunks")

    except Exception as e:
        logger.exception(f"Ingestion failed for document {doc_id}: {e}")
        await db.rollback()
        # Update status to failed
        try:
            doc_reload: KnowledgeDocument = await db.get(KnowledgeDocument, doc_id)
            if doc_reload:
                doc_reload.embedding_status = 3
                await db.commit()
        except Exception:
            pass
        raise


# ---------------------------------------------------------------------------
# Keyword search (PostgreSQL tsvector)
# ---------------------------------------------------------------------------

async def keyword_search(
    query: str,
    tenant_id: int,
    db: AsyncSession,
    top_k: int = 20,
) -> list[dict[str, Any]]:
    """
    Full-text search using PostgreSQL tsvector (Chinese + English).
    Falls back to ILIKE if no tsvector match.
    """
    # Build tsquery: split on whitespace, join with &
    words = query.strip().split()
    if not words:
        return []

    # Try plainto_tsquery first (handles Chinese tokenization via pg_jieba if available)
    try:
        sql = text(
            """
            SELECT
                kc.id AS chunk_id,
                kc.doc_id,
                kc.content,
                kc.heading,
                kc.page_number,
                kc.content_type,
                ts_rank_cd(to_tsvector('simple', kc.content), plainto_tsquery('simple', :query)) AS rank
            FROM knowledge_chunk kc
            JOIN knowledge_document kd ON kd.id = kc.doc_id
            WHERE kd.tenant_id = :tenant_id
              AND kd.embedding_status = 2
              AND to_tsvector('simple', kc.content) @@ plainto_tsquery('simple', :query)
            ORDER BY rank DESC
            LIMIT :top_k
            """
        )
        result = await db.execute(sql, {"query": query, "tenant_id": tenant_id, "top_k": top_k})
        rows = result.fetchall()

        if rows:
            return [
                {
                    "chunk_id": r.chunk_id,
                    "doc_id": r.doc_id,
                    "score": float(r.rank),
                    "content": r.content,
                    "heading": r.heading,
                    "page_number": r.page_number,
                    "content_type": r.content_type,
                }
                for r in rows
            ]
    except Exception as e:
        logger.warning(f"tsvector search failed, falling back to ILIKE: {e}")

    # Fallback: ILIKE
    ilike_sql = text(
        """
        SELECT
            kc.id AS chunk_id,
            kc.doc_id,
            kc.content,
            kc.heading,
            kc.page_number,
            kc.content_type,
            1.0 AS rank
        FROM knowledge_chunk kc
        JOIN knowledge_document kd ON kd.id = kc.doc_id
        WHERE kd.tenant_id = :tenant_id
          AND kd.embedding_status = 2
          AND kc.content ILIKE :pattern
        LIMIT :top_k
        """
    )
    pattern = f"%{query}%"
    result = await db.execute(ilike_sql, {"tenant_id": tenant_id, "pattern": pattern, "top_k": top_k})
    rows = result.fetchall()
    return [
        {
            "chunk_id": r.chunk_id,
            "doc_id": r.doc_id,
            "score": float(r.rank),
            "content": r.content,
            "heading": r.heading,
            "page_number": r.page_number,
            "content_type": r.content_type,
        }
        for r in rows
    ]


# ---------------------------------------------------------------------------
# RRF fusion
# ---------------------------------------------------------------------------

def _rrf_merge(
    vector_hits: list[dict],
    keyword_hits: list[dict],
    top_n: int = 10,
    k: int = 60,
) -> list[dict]:
    """
    Reciprocal Rank Fusion.
    Each hit gets score = 1/(k + rank) summed across lists.
    """
    scores: dict[int, float] = {}
    data: dict[int, dict] = {}

    for rank, hit in enumerate(vector_hits):
        cid = hit["chunk_id"]
        scores[cid] = scores.get(cid, 0.0) + 1.0 / (k + rank + 1)
        data[cid] = hit

    for rank, hit in enumerate(keyword_hits):
        cid = hit["chunk_id"]
        scores[cid] = scores.get(cid, 0.0) + 1.0 / (k + rank + 1)
        if cid not in data:
            data[cid] = hit

    ranked = sorted(scores.items(), key=lambda x: x[1], reverse=True)
    results = []
    for cid, rrf_score in ranked[:top_n]:
        entry = dict(data[cid])
        entry["rrf_score"] = rrf_score
        results.append(entry)

    return results


# ---------------------------------------------------------------------------
# Public hybrid search API
# ---------------------------------------------------------------------------

async def hybrid_search(
    query: str,
    tenant_id: int,
    db: AsyncSession,
    top_n: int = 8,
    doc_type: str | None = None,
    industry: str | None = None,
) -> list[dict[str, Any]]:
    """
    Hybrid search: vector similarity (Qdrant) + keyword (PostgreSQL) → RRF.
    Returns top_n chunks with rrf_score, content, heading, etc.
    """
    # Run both searches concurrently
    vector_task = asyncio.create_task(
        vector_search(query, tenant_id, top_k=20, doc_type=doc_type, industry=industry)
    )
    keyword_task = asyncio.create_task(
        keyword_search(query, tenant_id, db, top_k=20)
    )

    try:
        vector_hits, keyword_hits = await asyncio.gather(vector_task, keyword_task)
    except Exception as e:
        # If vector search fails (e.g. Qdrant not running), fall back to keyword only
        logger.warning(f"Vector search failed, using keyword only: {e}")
        keyword_hits = await keyword_search(query, tenant_id, db, top_k=20)
        vector_hits = []

    return _rrf_merge(vector_hits, keyword_hits, top_n=top_n)


# ---------------------------------------------------------------------------
# Delete document + vectors
# ---------------------------------------------------------------------------

async def delete_document(doc_id: int, db: AsyncSession) -> None:
    """Delete document, all its chunks (cascade), and Qdrant vectors."""
    from app.services.embedding_service import delete_document_vectors

    # Delete from Qdrant first (best-effort)
    try:
        await delete_document_vectors(doc_id)
    except Exception as e:
        logger.warning(f"Failed to delete Qdrant vectors for doc {doc_id}: {e}")

    # Delete from DB (chunks cascade via FK)
    doc = await db.get(KnowledgeDocument, doc_id)
    if doc:
        await db.delete(doc)
        await db.commit()
