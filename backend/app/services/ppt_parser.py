import io
import subprocess
import tempfile
import os
from dataclasses import dataclass, field
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from PIL import Image


@dataclass
class ParsedPage:
    page_number: int
    title: str
    content: str  # extracted text body
    speaker_notes: str
    thumbnail_bytes: bytes  # PNG


@dataclass
class ParsedPPT:
    page_count: int
    pages: list[ParsedPage] = field(default_factory=list)


def parse_pptx(file_bytes: bytes) -> ParsedPPT:
    """Parse .pptx or .pdf: extract text per page + render thumbnails."""
    if file_bytes[:4] == b"%PDF":
        return _parse_pdf(file_bytes)
    return _parse_pptx_bytes(file_bytes)


def _parse_pptx_bytes(file_bytes: bytes) -> ParsedPPT:
    with tempfile.TemporaryDirectory() as tmp_dir:
        pptx_path = os.path.join(tmp_dir, "deck.pptx")
        with open(pptx_path, "wb") as f:
            f.write(file_bytes)

        prs = Presentation(io.BytesIO(file_bytes))
        pages: list[ParsedPage] = []
        thumbnails = _render_thumbnails(pptx_path, tmp_dir, len(prs.slides))

        for idx, slide in enumerate(prs.slides):
            page_num = idx + 1
            title = _extract_title(slide)
            content = _extract_body(slide)
            notes = _extract_notes(slide)
            pages.append(ParsedPage(
                page_number=page_num,
                title=title,
                content=content,
                speaker_notes=notes,
                thumbnail_bytes=thumbnails.get(page_num, b""),
            ))

        return ParsedPPT(page_count=len(pages), pages=pages)


def _parse_pdf(file_bytes: bytes) -> ParsedPPT:
    """Parse a PDF using pdftotext (text) + pdftoppm (thumbnails)."""
    with tempfile.TemporaryDirectory() as tmp_dir:
        pdf_path = os.path.join(tmp_dir, "deck.pdf")
        with open(pdf_path, "wb") as f:
            f.write(file_bytes)

        # Get page count via pdfinfo
        page_count = _pdf_page_count(pdf_path)
        if page_count == 0:
            raise ValueError("PDF has no readable pages")

        # Thumbnails: pdftoppm directly on PDF
        png_prefix = os.path.join(tmp_dir, "slide")
        try:
            subprocess.run(
                ["pdftoppm", "-png", "-r", "96", pdf_path, png_prefix],
                capture_output=True, timeout=60, check=True,
            )
        except (subprocess.CalledProcessError, FileNotFoundError):
            pass

        pages: list[ParsedPage] = []
        for i in range(1, page_count + 1):
            # Extract text for this page
            text = _pdf_page_text(pdf_path, i)
            lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
            title = lines[0] if lines else ""
            content = "\n".join(lines[1:])

            # Find thumbnail
            thumbnail = b""
            for pad in (1, 2, 3):
                path = f"{png_prefix}-{i:0{pad}d}.png"
                if os.path.exists(path):
                    thumbnail = _resize_png(path, 640, 480)
                    break

            pages.append(ParsedPage(
                page_number=i,
                title=title,
                content=content,
                speaker_notes="",
                thumbnail_bytes=thumbnail,
            ))

        return ParsedPPT(page_count=page_count, pages=pages)


def _pdf_page_count(pdf_path: str) -> int:
    try:
        result = subprocess.run(
            ["pdfinfo", pdf_path],
            capture_output=True, text=True, timeout=10,
        )
        for line in result.stdout.splitlines():
            if line.lower().startswith("pages:"):
                return int(line.split(":")[1].strip())
    except Exception:
        pass
    return 0


def _pdf_page_text(pdf_path: str, page: int) -> str:
    try:
        result = subprocess.run(
            ["pdftotext", "-f", str(page), "-l", str(page), "-layout", pdf_path, "-"],
            capture_output=True, text=True, timeout=15,
        )
        return result.stdout
    except Exception:
        return ""


def _extract_title(slide) -> str:
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        return slide.shapes.title.text_frame.text.strip()
    return ""


def _extract_body(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    parts.append(text)
        if shape.has_table:
            for row in shape.table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip(" |"):
                    parts.append(row_text)
    return "\n".join(parts)


def _extract_notes(slide) -> str:
    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame:
            return notes_frame.text.strip()
    return ""


_SOFFICE = next(
    (p for p in ["/usr/local/bin/soffice", "/usr/bin/soffice", "soffice", "libreoffice"]
     if subprocess.run(["which", p.split("/")[-1]], capture_output=True).returncode == 0
     or os.path.exists(p)),
    "soffice",
)


def _render_thumbnails(pptx_path: str, output_dir: str, page_count: int) -> dict[int, bytes]:
    """Convert PPTX → PDF via LibreOffice, then PDF → PNG per page via pdftoppm."""
    thumbnails: dict[int, bytes] = {}

    # Step 1: PPTX → PDF
    pdf_dir = os.path.join(output_dir, "pdf")
    os.makedirs(pdf_dir, exist_ok=True)
    try:
        subprocess.run(
            [_SOFFICE, "--headless", "--convert-to", "pdf", "--outdir", pdf_dir, pptx_path],
            capture_output=True, timeout=60, check=True,
        )
    except (subprocess.CalledProcessError, FileNotFoundError):
        return thumbnails

    stem = os.path.splitext(os.path.basename(pptx_path))[0]
    pdf_path = os.path.join(pdf_dir, f"{stem}.pdf")
    if not os.path.exists(pdf_path):
        return thumbnails

    # Step 2: PDF → PNG per page via pdftoppm
    png_prefix = os.path.join(output_dir, "slide")
    try:
        subprocess.run(
            ["pdftoppm", "-png", "-r", "96", pdf_path, png_prefix],
            capture_output=True, timeout=60, check=True,
        )
    except (subprocess.CalledProcessError, FileNotFoundError):
        return thumbnails

    # pdftoppm outputs slide-1.png, slide-2.png ... or slide-01.png etc.
    for i in range(1, page_count + 1):
        for pad in (1, 2, 3):
            path = f"{png_prefix}-{i:0{pad}d}.png"
            if os.path.exists(path):
                thumbnails[i] = _resize_png(path, 640, 480)
                break

    return thumbnails


def _resize_png(path: str, width: int, height: int) -> bytes:
    with Image.open(path) as img:
        img.thumbnail((width, height), Image.LANCZOS)
        buf = io.BytesIO()
        img.save(buf, format="PNG", optimize=True)
        return buf.getvalue()
